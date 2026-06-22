"""
Generate KIA PhotoBooth operator documentation as a PowerPoint file.
Run: python scripts/generate-photobooth-docs-pptx.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "KIA-PhotoBooth-Documentation.pptx"

# KIA-inspired palette
KIA_RED = RGBColor(187, 22, 26)
DARK = RGBColor(28, 28, 30)
MID = RGBColor(90, 90, 95)
LIGHT_BG = RGBColor(245, 245, 247)
WHITE = RGBColor(255, 255, 255)
PLACEHOLDER_FILL = RGBColor(220, 222, 228)
PLACEHOLDER_BORDER = RGBColor(160, 164, 172)


def set_run_font(run, size=18, bold=False, color=DARK):
    run.font.name = "Segoe UI"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_title_slide(prs: Presentation, title: str, subtitle: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK
    bg.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.5), Inches(2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    set_run_font(r, size=40, bold=True, color=WHITE)

    sub = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(11.5), Inches(1.2))
    stf = sub.text_frame
    sp = stf.paragraphs[0]
    sr = sp.add_run()
    sr.text = subtitle
    set_run_font(sr, size=20, color=RGBColor(200, 200, 205))

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.8), Inches(3.9), Inches(2.2), Inches(0.06)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = KIA_RED
    accent.line.fill.background()


def add_section_slide(prs: Presentation, section: str, description: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = KIA_RED
    bg.line.fill.background()

    tbox = slide.shapes.add_textbox(Inches(0.9), Inches(2.8), Inches(11), Inches(1.5))
    tf = tbox.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = section
    set_run_font(r, size=36, bold=True, color=WHITE)

    if description:
        dbox = slide.shapes.add_textbox(Inches(0.9), Inches(4.0), Inches(11), Inches(1))
        dtf = dbox.text_frame
        dp = dtf.paragraphs[0]
        dr = dp.add_run()
        dr.text = description
        set_run_font(dr, size=18, color=WHITE)


def add_image_placeholder(slide, left, top, width, height, label: str):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PLACEHOLDER_FILL
    shape.line.color.rgb = PLACEHOLDER_BORDER
    shape.line.width = Pt(1.5)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    set_run_font(r, size=14, bold=True, color=MID)


def add_bullets(text_frame, lines: list[str], size=15, color=DARK, spacing=6):
    for i, line in enumerate(lines):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.text = line
        p.level = 0
        p.space_after = Pt(spacing)
        p.font.name = "Segoe UI"
        p.font.size = Pt(size)
        p.font.color.rgb = color


def add_content_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
    placeholder_label: str,
    footer: str = "",
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # header bar
    header = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = DARK
    header.line.fill.background()

    hbox = slide.shapes.add_textbox(Inches(0.55), Inches(0.18), Inches(12), Inches(0.6))
    hp = hbox.text_frame.paragraphs[0]
    hr = hp.add_run()
    hr.text = title
    set_run_font(hr, size=24, bold=True, color=WHITE)

    add_image_placeholder(
        slide,
        Inches(0.55),
        Inches(1.15),
        Inches(6.1),
        Inches(5.35),
        placeholder_label,
    )

    bbox = slide.shapes.add_textbox(Inches(6.95), Inches(1.15), Inches(5.9), Inches(5.35))
    btf = bbox.text_frame
    btf.word_wrap = True
    add_bullets(btf, bullets, size=14)

    if footer:
        fbox = slide.shapes.add_textbox(Inches(0.55), Inches(6.65), Inches(12.3), Inches(0.45))
        fp = fbox.text_frame.paragraphs[0]
        fr = fp.add_run()
        fr.text = footer
        set_run_font(fr, size=11, color=MID)


def add_flow_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = DARK
    header.line.fill.background()

    hbox = slide.shapes.add_textbox(Inches(0.55), Inches(0.18), Inches(12), Inches(0.6))
    hr = hbox.text_frame.paragraphs[0].add_run()
    hr.text = "End-to-end guest flow (the whole journey)"
    set_run_font(hr, size=24, bold=True, color=WHITE)

    steps = [
        ("BEFORE BOOTH", "Guest registers at KIA hub / event app and receives a unique QR code on their phone."),
        ("1 · SCAN", "Guest scans QR at the booth (hardware scanner, booth camera, or typed code)."),
        ("2 · CAPTURE", "Booth opens camera, plays countdown timer, takes photo automatically."),
        ("3 · RESULT", "Guest previews photo, picks an optional frame, taps save."),
        ("4 · UPLOAD", "Photo uploads to KIA hub gallery (works offline — syncs later)."),
        ("5 · HOME", "Success screen → booth returns to scan screen for next guest."),
    ]

    y = Inches(1.2)
    for i, (label, text) in enumerate(steps):
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(0.55),
            y,
            Inches(12.3),
            Inches(0.82),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BG if i % 2 == 0 else WHITE
        box.line.color.rgb = PLACEHOLDER_BORDER

        lbox = slide.shapes.add_textbox(Inches(0.75), y + Inches(0.12), Inches(2.2), Inches(0.5))
        lp = lbox.text_frame.paragraphs[0]
        lr = lp.add_run()
        lr.text = label
        set_run_font(lr, size=13, bold=True, color=KIA_RED)

        tbox = slide.shapes.add_textbox(Inches(2.9), y + Inches(0.1), Inches(9.6), Inches(0.6))
        tp = tbox.text_frame.paragraphs[0]
        tr = tp.add_run()
        tr.text = text
        set_run_font(tr, size=13, color=DARK)

        y += Inches(0.92)

    note = slide.shapes.add_textbox(Inches(0.55), Inches(6.55), Inches(12.3), Inches(0.5))
    nr = note.text_frame.paragraphs[0].add_run()
    nr.text = "Tip: One QR scan = one session. After upload, the guest finds their photo in the KIA hub photogallery."
    set_run_font(nr, size=11, color=MID)


def add_two_column_slide(prs: Presentation, title: str, left_title: str, left_bullets: list[str], right_title: str, right_bullets: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = DARK
    header.line.fill.background()

    hbox = slide.shapes.add_textbox(Inches(0.55), Inches(0.18), Inches(12), Inches(0.6))
    hr = hbox.text_frame.paragraphs[0].add_run()
    hr.text = title
    set_run_font(hr, size=24, bold=True, color=WHITE)

    for col, (ctitle, cbullets, x) in enumerate(
        [(left_title, left_bullets, 0.55), (right_title, right_bullets, 6.75)]
    ):
        cbox = slide.shapes.add_textbox(Inches(x), Inches(1.1), Inches(5.9), Inches(0.4))
        cr = cbox.text_frame.paragraphs[0].add_run()
        cr.text = ctitle
        set_run_font(cr, size=16, bold=True, color=KIA_RED)

        body = slide.shapes.add_textbox(Inches(x), Inches(1.55), Inches(5.9), Inches(5.2))
        add_bullets(body.text_frame, cbullets, size=13)


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "KIA PhotoBooth",
        "Simple operator guide — every screen, every step, plus admin",
    )

    add_section_slide(prs, "What is this?", "A touchscreen kiosk app for event photo keepsakes")

    add_content_slide(
        prs,
        "What the booth does",
        [
            "Runs full-screen on a Windows PC at the event.",
            "Guests prove they registered by scanning a personal QR code.",
            "A DSLR (or webcam fallback) takes their photo after a countdown.",
            "Guests pick an optional branded frame overlay.",
            "The finished photo uploads to the KIA hub photogallery.",
            "The booth resets automatically for the next guest.",
            "",
            "App routes:",
            "  /          → QR scan (home)",
            "  /capture   → Camera & countdown",
            "  /result    → Preview, frames, upload",
            "  /admin     → Operator settings (PIN required)",
        ],
        "[INSERT IMAGE: Booth hardware setup — screen, camera, optional QR scanner]",
    )

    add_flow_slide(prs)

    add_section_slide(prs, "Guest screens", "Four screens — scan, pose, review, done")

    add_content_slide(
        prs,
        "Screen 1 — QR Scan (Home screen)",
        [
            "Route: /  (this is the idle / waiting screen)",
            "",
            "What the guest sees:",
            "• KIA logo + tagline (“Movement that inspires”)",
            "• Headline: “Capture Your KIA Moment”",
            "• Animated QR scan area with moving scan line",
            "• Status text under the QR frame",
            "",
            "What the guest does:",
            "• Hold their registration QR up to the scanner, OR",
            "• Tap the QR area to use the booth camera (if enabled), OR",
            "• Operator can type a code on a USB keyboard + Enter",
            "",
            "What happens:",
            "• Code is validated against KIA API",
            "• Green checkmark + “QR verified”",
            "• After ~1 second → Capture screen",
            "",
            "If invalid: error message stays on screen; guest tries again.",
        ],
        "[INSERT SCREENSHOT: QR scan screen — idle state]",
        footer="Hidden Admin link: hold the logo for 6 seconds to reveal it for 20 minutes.",
    )

    add_content_slide(
        prs,
        "Screen 1 — QR Scan (states)",
        [
            "Status messages (under the QR frame):",
            "",
            "Idle:     “Scan your QR code to get started”",
            "Scanning: “Scanning…” (camera mode)",
            "Checking: “Verifying…” (API call in progress)",
            "Success:  “QR verified” + animated checkmark",
            "",
            "Scanner options (configured in Admin):",
            "• GFS4400 USB serial scanner (default at events)",
            "• Booth webcam reads QR from guest’s phone",
            "• Manual keyboard entry for testing",
            "",
            "Testing bypass:",
            "• Configured bypass code (default: 12345) skips normal QR format",
            "• Useful for setup without real guest codes",
        ],
        "[INSERT SCREENSHOT: QR scan — success checkmark overlay]",
    )

    add_content_slide(
        prs,
        "Screen 2 — Capture (Get Ready)",
        [
            "Route: /capture  (opens automatically after QR success)",
            "",
            "What the guest sees:",
            "• “Get Ready” / “Strike Your Best Pose”",
            "• Live camera preview in a framed window",
            "• Countdown timer video on the right",
            "• Footer: “Look at camera and hold still.”",
            "",
            "What the guest does:",
            "• Nothing — the booth runs automatically",
            "• Pose and stay still during the countdown",
            "",
            "What happens (automatic):",
            "1. Camera initializes (Canon SDK or webcam)",
            "2. Countdown timer video plays (~6–7 seconds)",
            "3. Photo is taken at the end of the timer",
            "4. App navigates to Result screen",
            "",
            "Camera note: portrait mount rotates preview 90° before crop.",
        ],
        "[INSERT SCREENSHOT: Capture screen — live preview + timer]",
        footer="No tap-to-capture — timer always drives the shot.",
    )

    add_content_slide(
        prs,
        "Screen 3 — Result (Your keepsake)",
        [
            "Route: /result  (shows captured photo)",
            "",
            "What the guest sees:",
            "• “Your KIA Keepsake is Ready”",
            "• Large preview of their captured photo",
            "• Row of frame buttons below the photo",
            "• Green tick button (save & upload)",
            "",
            "What the guest does:",
            "1. Tap a frame to preview overlay (optional)",
            "   — “No effect” = plain photo",
            "   — Some frames may show a lock icon (not available)",
            "2. Tap the green ✓ to confirm",
            "",
            "What happens:",
            "• Upload is queued immediately (offline-safe)",
            "• “Uploading your photo…” shows for min. 3 seconds",
            "• Then → Upload Success screen",
        ],
        "[INSERT SCREENSHOT: Result screen — photo + frame picker]",
    )

    add_content_slide(
        prs,
        "Screen 4 — Upload success",
        [
            "Shown after guest taps the green save button.",
            "",
            "What the guest sees:",
            "• “Upload completed successfully”",
            "• “Check your KIA hub photogallery for the picture”",
            "• Optional countdown: “Returning to home in Xs”",
            "• “Continue” button",
            "",
            "What the guest does:",
            "• Check their phone / KIA hub for the gallery photo",
            "• Tap Continue (or wait for auto-return)",
            "",
            "What happens:",
            "• Booth clears the session",
            "• Camera session closes",
            "• Returns to QR Scan screen for next guest",
            "• Auto-home default: 10 seconds (configurable in Admin)",
            "",
            "Upload runs in background — booth does not wait for network.",
        ],
        "[INSERT SCREENSHOT: Upload success screen]",
    )

    add_section_slide(prs, "Admin panel", "For operators & tech setup — not shown to guests")

    add_content_slide(
        prs,
        "How to open Admin",
        [
            "Admin is hidden from guests on purpose.",
            "",
            "Step 1 — Reveal the link:",
            "• On the QR scan screen, press and hold the KIA logo",
            "• Hold for 6 full seconds",
            "• “Admin” link appears bottom-left for 20 minutes",
            "",
            "Step 2 — Open Admin:",
            "• Tap “Admin” → goes to /admin/login",
            "",
            "Step 3 — Sign in:",
            "• Enter operator PIN on the touch numpad",
            "• Default shipped PIN: 2727",
            "  (check photobooth-config.json if changed)",
            "",
            "Also from login screen:",
            "• “Back” → return to booth",
            "• “Shutdown system” → power off PC (end of event)",
        ],
        "[INSERT SCREENSHOT: Admin link visible on QR screen]",
        footer="Change PIN in config/photobooth-config.json → adminPin",
    )

    add_content_slide(
        prs,
        "Admin — Login screen",
        [
            "Route: /admin/login",
            "",
            "Features:",
            "• Touch-friendly PIN entry (on-screen numpad)",
            "• ⓘ button shows default PIN hint",
            "• Sign in → Admin dashboard",
            "• Back → guest booth",
            "",
            "Shutdown system (red button):",
            "• Confirmation dialog before action",
            "• Closes PhotoBooth app AND powers off the PC",
            "• Use only when the event is finished",
            "",
            "Security:",
            "• Dashboard requires valid PIN (adminGuard)",
            "• Session ends when you tap Sign out",
        ],
        "[INSERT SCREENSHOT: Admin login + numpad]",
    )

    add_content_slide(
        prs,
        "Admin — Dashboard overview",
        [
            "Route: /admin  (after PIN login)",
            "",
            "Top bar buttons:",
            "• Reload — refresh config from disk",
            "• Sign out — lock admin",
            "• Booth — jump back to guest screen",
            "",
            "Three tabs in the sidebar:",
            "1. Copy      — all on-screen text",
            "2. Logo      — branding image & size",
            "3. Scanner & API — hardware + backend",
            "",
            "Changes save to photobooth-config.json",
            "Most changes apply live; some need Reload.",
        ],
        "[INSERT SCREENSHOT: Admin dashboard — full view]",
    )

    add_content_slide(
        prs,
        "Admin tab — Copy",
        [
            "Edit every string guests see on the booth.",
            "",
            "Scan landing (QR screen):",
            "Tagline, title, subtitle, scan prompts, error messages, admin link label",
            "",
            "Capture screen:",
            "Starting camera text, ready title/subtitle, footer hint",
            "",
            "Result & upload:",
            "Keepsake title, upload messages, success text,",
            "upload min display seconds (default 3),",
            "auto-return home seconds (default 10, 0 = off)",
            "",
            "Click “Save copy” at the bottom when done.",
        ],
        "[INSERT SCREENSHOT: Admin — Copy tab]",
    )

    add_content_slide(
        prs,
        "Admin tab — Logo (Branding)",
        [
            "Control the logo shown on all guest screens.",
            "",
            "Options:",
            "• Upload custom PNG, JPG, SVG, or WebP",
            "• Adjust size slider: 50% – 200%",
            "• “Use default KIA logo” to reset",
            "• Live preview in the admin panel",
            "",
            "Reminder shown in admin:",
            "Admin link is hidden until logo is held 6 seconds.",
            "",
            "Click “Save logo size” after adjusting scale.",
            "Use “Choose logo image…” to pick a new file.",
        ],
        "[INSERT SCREENSHOT: Admin — Logo tab with preview]",
    )

    add_two_column_slide(
        prs,
        "Admin tab — Scanner & API (part 1)",
        "Booth camera",
        [
            "Camera orientation:",
            "• Portrait (default) — rotates 90° for side-mounted DSLR",
            "• Landscape — sensor already upright",
            "",
            "GFS4400 serial scanner:",
            "• Enable/disable USB-COM scanner",
            "• Pick COM port from dropdown",
            "• Baud rate: 9600 (standard)",
            "• “Test open port” to verify connection",
            "• Status + last scanned code shown live",
            "",
            "Fallback:",
            "☑ Use booth camera for QR when scanner unavailable",
        ],
        "KIA Forum API",
        [
            "API base URL:",
            "• Dev: dev-kiaforum2026.thetunagroup.com",
            "• Prod: kiaexperience.info",
            "",
            "Photo upload base URL (optional):",
            "• Separate host for media upload only",
            "",
            "Other settings:",
            "• Bearer token (optional bootstrap)",
            "• Dev bypass email",
            "• QR token prefix (KIA-PHOTO-)",
            "• Bypass code for testing",
            "• Offline prefix validation toggle",
            "• API debug log panel toggle",
            "• Upload format: PNG or JPEG",
            "",
            "“Test API connection” + upload queue status",
        ],
    )

    add_content_slide(
        prs,
        "Admin tab — Scanner & API (part 2)",
        [
            "Save buttons (important — there are two):",
            "• “Save scanner & camera settings”",
            "• “Save API settings”",
            "",
            "Typical event-day checklist:",
            "1. Refresh ports → pick correct COM port",
            "2. Test open port → should succeed",
            "3. Test API connection → should succeed",
            "4. Scan a real guest QR on home screen",
            "5. Take a test photo end-to-end",
            "6. Confirm photo appears in KIA hub gallery",
            "",
            "Offline mode:",
            "If API is down but prefix validation is on,",
            "guests can still scan and photos queue for later sync.",
        ],
        "[INSERT SCREENSHOT: Admin — Scanner & API tab]",
    )

    add_section_slide(prs, "Quick reference", "Cheat sheet for operators")

    add_two_column_slide(
        prs,
        "Operator cheat sheet",
        "Guest stuck? Try this",
        [
            "Stuck on QR screen → rescan; check scanner cable",
            "“Invalid QR code” → guest needs event registration QR",
            "Camera black → check DSLR USB + power; Reload config",
            "Upload error on result → tap ✓ again; check API in Admin",
            "Booth frozen → Alt+F4 app or restart PC",
            "End of event → Admin login → Shutdown system",
        ],
        "Default values (shipped config)",
        [
            "Admin PIN: 2727",
            "Bypass test code: 12345",
            "QR prefix: KIA-PHOTO-",
            "Scanner baud: 9600",
            "Upload min screen: 3 seconds",
            "Auto home: 10 seconds",
            "Admin link hold: 6 sec reveal / 20 min visible",
            "Config file: config/photobooth-config.json",
        ],
    )

    add_content_slide(
        prs,
        "Image placeholder checklist",
        [
            "Replace each gray box with a real screenshot before sharing:",
            "",
            "☐ Booth hardware setup photo",
            "☐ QR scan — idle",
            "☐ QR scan — success checkmark",
            "☐ Capture — live preview + timer",
            "☐ Result — photo + frames",
            "☐ Upload success screen",
            "☐ Admin link visible on QR screen",
            "☐ Admin login + numpad",
            "☐ Admin dashboard overview",
            "☐ Admin Copy tab",
            "☐ Admin Logo tab",
            "☐ Admin Scanner & API tab",
            "",
            "Tip: run the app fullscreen, use Win+Shift+S to capture.",
        ],
        "[INSERT IMAGE: Example annotated screenshot with labels]",
    )

    add_title_slide(
        prs,
        "Questions?",
        "KIA PhotoBooth — Movement that inspires",
    )

    return prs


def main():
    prs = build_presentation()
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Created: {OUTPUT}")


if __name__ == "__main__":
    main()
